"""Create deterministic G01 golden and source-security fixtures."""

from __future__ import annotations

import argparse
import json
import zipfile
from io import BytesIO
from pathlib import Path

from instant_ppt_worker.paths import REPOSITORY_ROOT
from instant_ppt_worker.security import scan_source
from instant_ppt_worker.source_parser import deterministic_ulid, parse_source
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

GOLDEN_ROOT = REPOSITORY_ROOT / "tests" / "golden"
THREAT_ROOT = REPOSITORY_ROOT / "tests" / "security-fixtures"
FIXED_ZIP_TIME = (2026, 8, 16, 0, 0, 0)
ORGANIZATION_ID = "01ARZ3NDEKTSV4RRFFQ69G5FAV"
CREATED_AT = "2026-08-16T00:00:00Z"

CASES = [
    {
        "slug": "01-zh-overview",
        "extension": ".md",
        "format": "markdown",
        "title": "年度经营复盘",
        "source": "# 年度经营复盘\n\n收入同比增长 18%。\n\n- 核心产品增长\n- 客户留存改善\n",
        "body": ["收入同比增长 18%", "核心产品贡献主要增量", "客户留存率持续改善"],
        "coverage": ["zh-CN", "overview"],
    },
    {
        "slug": "02-en-market",
        "extension": ".txt",
        "format": "text",
        "title": "Market Expansion Brief",
        "source": (
            "Market Expansion Brief\nNorth America grew 24 percent.\n"
            "Retention reached 91 percent.\n"
        ),
        "body": ["North America grew 24%", "Retention reached 91%", "Enterprise leads increased"],
        "coverage": ["en-US", "plain-text"],
    },
    {
        "slug": "03-long-title",
        "extension": ".html",
        "format": "html",
        "title": "面向跨区域复杂业务协同场景的下一代智能运营平台阶段性成果总结",
        "source": (
            "<html><body><h1>长标题压力测试</h1><p>信息应保持清晰且可编辑。</p></body></html>"
        ),
        "body": ["长标题自动采用紧凑字号", "信息层级保持清晰", "输出仍为原生可编辑文本"],
        "coverage": ["long-title", "html"],
    },
    {
        "slug": "04-table-docx",
        "extension": ".docx",
        "format": "docx",
        "title": "区域经营数据表",
        "source": "区域|收入|增长\n华东|3200|18%\n华南|2400|15%",
        "body": ["华东收入 3,200 万", "华南收入 2,400 万", "表格结构完成解析"],
        "coverage": ["table", "docx"],
    },
    {
        "slug": "05-chart-pptx",
        "extension": ".pptx",
        "format": "pptx",
        "title": "季度增长趋势",
        "source": "Q1 12%\nQ2 16%\nQ3 21%\nQ4 27%",
        "body": ["Q1 增长 12%", "Q2 增长 16%", "Q4 达到 27%"],
        "coverage": ["chart", "pptx-source"],
    },
    {
        "slug": "06-mixed-fonts",
        "extension": ".md",
        "format": "markdown",
        "title": "字体回退 Font Fallback",
        "source": "# 字体回退 Font Fallback\n\n中文、English 与数字 2026 混排。\n",
        "body": ["Microsoft YaHei 中文回退", "Arial Latin fallback", "数字与符号保持可编辑"],
        "coverage": ["mixed-language", "font-fallback"],
    },
    {
        "slug": "07-template-brand",
        "extension": ".html",
        "format": "html",
        "title": "品牌模板一致性",
        "source": (
            "<html><body><h1>Brand System</h1><p>Blue, slate and white tokens.</p></body></html>"
        ),
        "body": ["主色使用品牌蓝", "背景使用中性灰", "模板绑定版本固定"],
        "coverage": ["template", "brand"],
    },
    {
        "slug": "08-dense-docx",
        "extension": ".docx",
        "format": "docx",
        "title": "密集要点信息压测",
        "source": "目标一\n目标二\n目标三\n风险一\n风险二\n行动计划",
        "body": [
            "目标一：稳定交付",
            "目标二：提升效率",
            "目标三：控制成本",
            "风险一：资源波动",
            "风险二：依赖变化",
            "行动：周度复盘",
        ],
        "coverage": ["dense-content", "six-bullets"],
    },
    {
        "slug": "09-pdf-baseline",
        "extension": ".pdf",
        "format": "pdf",
        "title": "Permissive PDF Baseline",
        "source": "PDF baseline uses pypdf without PyMuPDF.",
        "body": [
            "Permissive parser: pypdf",
            "Encrypted files are rejected",
            "EPUB remains disabled",
        ],
        "coverage": ["pdf", "license-boundary"],
    },
    {
        "slug": "10-multilingual-pptx",
        "extension": ".pptx",
        "format": "pptx",
        "title": "全球发布 Global Launch",
        "source": "全球发布\nGlobal Launch\nTokyo Shanghai New York",
        "body": ["上海：产品发布", "Tokyo: partner event", "New York: customer forum"],
        "coverage": ["multilingual", "pptx-source"],
    },
]


def _write_json(path: Path, value: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def _reproducible_zip(path: Path) -> None:
    with zipfile.ZipFile(path) as source:
        members = [
            (info.filename, source.read(info), info.compress_type) for info in source.infolist()
        ]
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w") as target:
        for name, data, compression in sorted(members):
            info = zipfile.ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = compression
            info.external_attr = 0o100644 << 16
            target.writestr(info, data)
    path.write_bytes(buffer.getvalue())


def _write_docx(path: Path, text: str) -> None:
    paragraphs = "".join(f"<w:p><w:r><w:t>{line}</w:t></w:r></w:p>" for line in text.splitlines())
    entries = {
        "[Content_Types].xml": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" '
            'ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            "</Types>"
        ),
        "_rels/.rels": (
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/'
            'relationships/officeDocument" '
            'Target="word/document.xml"/></Relationships>'
        ),
        "word/document.xml": (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            f"<w:body>{paragraphs}<w:sectPr/></w:body></w:document>"
        ),
    }
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, data in entries.items():
            info = zipfile.ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o100644 << 16
            archive.writestr(info, data.encode("utf-8"))


def _write_pptx(path: Path, title: str, text: str) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = text
    presentation.save(path)
    _reproducible_zip(path)


def _write_pdf(path: Path, text: str) -> None:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 18 Tf 72 720 Td ({escaped}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
        ),
        f"<< /Length {len(stream)} >>\nstream\n".encode("ascii") + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, 1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode("ascii") + obj + b"\nendobj\n")
    xref = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode(
            "ascii"
        )
    )
    path.write_bytes(output)


def _source(case: dict[str, object], path: Path) -> None:
    format_name = case["format"]
    if format_name in {"markdown", "text", "html"}:
        path.write_text(str(case["source"]), encoding="utf-8")
    elif format_name == "docx":
        _write_docx(path, str(case["source"]))
    elif format_name == "pptx":
        _write_pptx(path, str(case["title"]), str(case["source"]))
    elif format_name == "pdf":
        _write_pdf(path, str(case["source"]))
    else:
        raise ValueError(f"unsupported fixture format: {format_name}")


def _deck_plan(case: dict[str, object], index: int) -> dict[str, object]:
    seed = f"{index:02x}" * 32
    snapshot_id = deterministic_ulid(seed)
    slides = []
    slide_content = [
        ("cover", str(case["title"]), [str(case["body"][0])]),
        ("content", "Key findings / 核心结论", list(case["body"])),
        (
            "summary",
            "Next step / 下一步",
            ["Validate editability", "Review visual baseline", "Record compatibility evidence"],
        ),
    ]
    for order, (role, title, body) in enumerate(slide_content):
        slide_id = deterministic_ulid(f"{index:02x}{order:02x}" * 16)
        slides.append(
            {
                "schemaVersion": 1,
                "slideId": slide_id,
                "outlineSlideId": slide_id,
                "order": order,
                "role": role,
                "title": title,
                "body": body,
                "editable": True,
            }
        )
    return {
        "schemaVersion": 1,
        "snapshotId": snapshot_id,
        "title": case["title"],
        "modeId": "native",
        "templateBinding": {
            "schemaVersion": 1,
            "templateId": deterministic_ulid("ab" * 32),
            "templateVersionId": deterministic_ulid(f"{index:02x}" * 32),
            "compatibilityVersion": "ppt-master-v4.7.0",
            "roleBindings": {
                "cover": "layout-cover",
                "content": "layout-content",
                "summary": "layout-summary",
            },
        },
        "slides": slides,
    }


def materialize_golden(approve: bool) -> None:
    for index, case in enumerate(CASES, 1):
        case_root = GOLDEN_ROOT / str(case["slug"])
        source_dir = case_root / "source"
        source_dir.mkdir(parents=True, exist_ok=True)
        source_path = source_dir / f"input{case['extension']}"
        _source(case, source_path)
        source_key = source_path.relative_to(case_root).as_posix()
        source_id = deterministic_ulid(f"{index + 32:02x}" * 32)
        deck_plan = _deck_plan(case, index)
        _write_json(case_root / "deck-plan.approved.json", deck_plan)
        _write_json(
            case_root / "case.json",
            {
                "schemaVersion": 1,
                "slug": case["slug"],
                "sourceKey": source_key,
                "sourceId": source_id,
                "organizationId": ORGANIZATION_ID,
                "createdAt": CREATED_AT,
                "coverage": case["coverage"],
            },
        )
        if approve:
            decision = scan_source(source_key, source_path)
            if decision.decision != "clean":
                raise RuntimeError(f"golden source rejected: {case['slug']}: {decision.findings}")
            decision_path = case_root / "generated" / "security-decision.json"
            _write_json(decision_path, decision.model_dump(by_alias=True, mode="json"))
            parsed = parse_source(
                source_key,
                source_path,
                decision_path,
                case_root / "generated" / "parse",
                source_id=source_id,
                organization_id=ORGANIZATION_ID,
                created_at=CREATED_AT,
            )
            _write_json(case_root / "source-package.expected.json", parsed["sourcePackage"])


def _zip_fixture(path: Path, entries: dict[str, bytes], *, symlink: str | None = None) -> None:
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, data in entries.items():
            info = zipfile.ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = (0o120777 if name == symlink else 0o100644) << 16
            archive.writestr(info, data)


def materialize_threats() -> None:
    THREAT_ROOT.mkdir(parents=True, exist_ok=True)
    threats: list[dict[str, str]] = []

    def record(name: str, code: str) -> Path:
        threats.append({"sourceKey": name, "expectedCode": code})
        return THREAT_ROOT / name

    record("magic-mismatch.pdf", "SOURCE_EXTENSION_MAGIC_MISMATCH").write_text(
        "not a PDF", encoding="utf-8"
    )
    record("corrupt.docx", "SOURCE_ARCHIVE_CORRUPT").write_bytes(b"PK\x03\x04broken")
    record("corrupt.pdf", "SOURCE_PDF_CORRUPT").write_bytes(b"%PDF-1.4\nbroken")
    record("active.html", "SOURCE_HTML_ACTIVE_CONTENT").write_text(
        "<html><script>alert(1)</script></html>", encoding="utf-8"
    )
    record("external.html", "SOURCE_HTML_EXTERNAL_REFERENCE").write_text(
        '<html><img src="https://example.invalid/pixel.png"></html>', encoding="utf-8"
    )
    record("virus-canary.txt", "SOURCE_MALWARE_TEST_SIGNATURE").write_text(
        "INSTANT-PPT-EICAR-TEST-SIGNATURE", encoding="utf-8"
    )
    encrypted = PdfWriter()
    encrypted.add_blank_page(width=612, height=792)
    encrypted.encrypt("fixture-password")
    with record("encrypted.pdf", "SOURCE_PDF_ENCRYPTED").open("wb") as stream:
        encrypted.write(stream)
    _zip_fixture(
        record("traversal.docx", "SOURCE_ZIP_PATH_TRAVERSAL"),
        {
            "[Content_Types].xml": b"<Types/>",
            "word/document.xml": b"<document/>",
            "../escape.txt": b"escape",
        },
    )
    deep = "/".join(["deep"] * 10) + "/document.xml"
    _zip_fixture(
        record("deep.docx", "SOURCE_ZIP_DEPTH_LIMIT"),
        {"[Content_Types].xml": b"<Types/>", "word/document.xml": b"<document/>", deep: b"x"},
    )
    _zip_fixture(
        record("ratio.docx", "SOURCE_ZIP_RATIO_LIMIT"),
        {"[Content_Types].xml": b"<Types/>", "word/document.xml": b"A" * (2 * 1024 * 1024)},
    )
    _zip_fixture(
        record("symlink.docx", "SOURCE_ZIP_SYMLINK"),
        {
            "[Content_Types].xml": b"<Types/>",
            "word/document.xml": b"<document/>",
            "word/link": b"target",
        },
        symlink="word/link",
    )
    rels = (
        b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        b'<Relationship Id="rId1" TargetMode="External" Target="https://example.invalid/payload"/>'
        b"</Relationships>"
    )
    _zip_fixture(
        record("external-rel.pptx", "SOURCE_OFFICE_EXTERNAL_RELATIONSHIP"),
        {
            "[Content_Types].xml": b"<Types/>",
            "ppt/presentation.xml": b"<presentation/>",
            "ppt/_rels/presentation.xml.rels": rels,
        },
    )
    many_entries = {"[Content_Types].xml": b"<Types/>", "word/document.xml": b"<document/>"}
    many_entries.update({f"word/items/{index:04d}.xml": b"x" for index in range(2_001)})
    _zip_fixture(record("entry-count.docx", "SOURCE_ZIP_ENTRY_LIMIT"), many_entries)
    _write_json(THREAT_ROOT / "manifest.json", {"schemaVersion": 1, "fixtures": threats})


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--approve", action="store_true", help="write expected SourcePackage files")
    arguments = parser.parse_args()
    materialize_golden(arguments.approve)
    materialize_threats()


if __name__ == "__main__":
    main()
