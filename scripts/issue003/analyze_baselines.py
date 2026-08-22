"""Generate deterministic ISSUE-003 PPTX and render-baseline evidence."""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import zipfile
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASELINE_FILES = {
    "reference-ppt-master": "reference-ppt-master.pptx",
    "before-deterministic-template": "before-deterministic-template.pptx",
    "before-user-download": "before-user-download.pptx",
}


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def shape_text(shape: Any) -> str:
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        return "".join(shape_text(child) for child in shape.shapes)
    if getattr(shape, "has_text_frame", False):
        return str(shape.text or "")
    return ""


def recursive_shape_count(shape: Any) -> int:
    if getattr(shape, "shape_type", None) != MSO_SHAPE_TYPE.GROUP:
        return 1
    return 1 + sum(recursive_shape_count(child) for child in shape.shapes)


def analyze_pptx(path: Path) -> dict[str, Any]:
    presentation = Presentation(path)
    used_layouts: set[str] = set()
    slides: list[dict[str, Any]] = []
    totals = {
        "topLevelShapeCount": 0,
        "recursiveShapeCount": 0,
        "visibleCharacterCount": 0,
        "pictureCount": 0,
        "nativeChartCount": 0,
        "nativeTableCount": 0,
    }
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [shape_text(shape) for shape in slide.shapes]
        visible_text = "".join(texts)
        top_level = len(slide.shapes)
        recursive = sum(recursive_shape_count(shape) for shape in slide.shapes)
        pictures = sum(
            1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )
        charts = sum(1 for shape in slide.shapes if getattr(shape, "has_chart", False))
        tables = sum(1 for shape in slide.shapes if getattr(shape, "has_table", False))
        used_layouts.add(str(slide.slide_layout.part.partname))
        record = {
            "slide": index,
            "topLevelShapeCount": top_level,
            "recursiveShapeCount": recursive,
            "visibleCharacterCount": len(visible_text),
            "pictureCount": pictures,
            "nativeChartCount": charts,
            "nativeTableCount": tables,
            "title": next((value.strip() for value in texts if value.strip()), "")[:240],
        }
        slides.append(record)
        for key in totals:
            totals[key] += int(record[key])
    with zipfile.ZipFile(path) as package:
        names = package.namelist()
        notes_pages = len(
            [
                name
                for name in names
                if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")
            ]
        )
        media_parts = len(
            [name for name in names if name.startswith("ppt/media/") and not name.endswith("/")]
        )
    return {
        "path": path.name,
        "sha256": sha256_file(path),
        "sizeBytes": path.stat().st_size,
        "slideCount": len(presentation.slides),
        "slideMasterCount": len(presentation.slide_masters),
        "usedSlideLayoutCount": len(used_layouts),
        "notesPageCount": notes_pages,
        "mediaPartCount": media_parts,
        **totals,
        "averageTopLevelShapesPerSlide": round(
            totals["topLevelShapeCount"] / max(1, len(presentation.slides)), 2
        ),
        "slides": slides,
    }


def natural_pngs(path: Path) -> list[Path]:
    def order(value: Path) -> tuple[int, str]:
        digits = "".join(character for character in value.stem if character.isdigit())
        return (int(digits) if digits else 0, value.name.casefold())

    unique = {
        candidate.resolve().as_posix().casefold(): candidate
        for candidate in path.iterdir()
        if candidate.is_file() and candidate.suffix.casefold() == ".png"
    }
    return sorted(unique.values(), key=order)


def contact_sheet(paths: list[Path], target: Path) -> dict[str, Any]:
    if not paths:
        return {"status": "missing", "path": target.name, "pngCount": 0}
    images = [Image.open(path).convert("RGB") for path in paths]
    thumb_width = 480
    thumb_height = round(thumb_width * 9 / 16)
    columns = 2
    rows = math.ceil(len(images) / columns)
    label_height = 34
    canvas = Image.new(
        "RGB",
        (columns * thumb_width, rows * (thumb_height + label_height)),
        "#111827",
    )
    draw = ImageDraw.Draw(canvas)
    for index, image in enumerate(images):
        thumb = image.copy()
        thumb.thumbnail((thumb_width, thumb_height), Image.Resampling.LANCZOS)
        x = (index % columns) * thumb_width
        y = (index // columns) * (thumb_height + label_height)
        canvas.paste(thumb, (x, y))
        draw.text((x + 12, y + thumb_height + 8), f"P{index + 1:02d}", fill="white")
    target.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(target, format="PNG", optimize=True)
    return {
        "status": "rendered",
        "path": target.name,
        "pngCount": len(paths),
        "sha256": sha256_file(target),
        "sizeBytes": target.stat().st_size,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("baseline_dir", type=Path)
    args = parser.parse_args()
    baseline_dir = args.baseline_dir.resolve()
    if not baseline_dir.is_dir():
        raise SystemExit(f"baseline directory does not exist: {baseline_dir}")
    metrics: dict[str, Any] = {
        "schemaVersion": 1,
        "comparisonPolicy": {
            "primaryBefore": "before-deterministic-template",
            "reference": "reference-ppt-master",
            "downloadLineageOnly": "before-user-download",
            "strictAB": False,
            "reason": (
                "The reference and website before artifacts have recoverable sources but not the "
                "same approved outline/model context. The frozen website snapshot is the authority "
                "for the implementation before/after rerun."
            ),
        },
        "decks": {},
        "renders": {},
        "frozenInputs": {},
    }
    for name, filename in BASELINE_FILES.items():
        path = baseline_dir / filename
        if path.is_file():
            metrics["decks"][name] = analyze_pptx(path)
        render_dir = baseline_dir / "renders" / name
        render_paths = natural_pngs(render_dir) if render_dir.is_dir() else []
        metrics["renders"][name] = contact_sheet(
            render_paths,
            baseline_dir / f"{name}-contact-sheet.png",
        )
    for filename in (
        "before-approved-snapshot.json",
        "before-approved-source.md",
        "before-conversion-profile.json",
        "reference-source-recovered.md",
        "reference-design-spec.md",
        "reference-spec-lock.md",
    ):
        path = baseline_dir / filename
        if path.is_file():
            metrics["frozenInputs"][filename] = {
                "sha256": sha256_file(path),
                "sizeBytes": path.stat().st_size,
            }
    output = baseline_dir / "baseline-metrics.json"
    output.write_text(
        json.dumps(metrics, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )
    print(
        json.dumps(
            {
                "output": str(output),
                "deckCount": len(metrics["decks"]),
                "renderedDeckCount": sum(
                    1 for value in metrics["renders"].values() if value["status"] == "rendered"
                ),
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
