"""PPTX package, relationship, editability, and content postflight."""

from __future__ import annotations

import json
import posixpath
import re
import zipfile
from collections import Counter
from pathlib import Path
from urllib.parse import unquote, urlsplit

from defusedxml import ElementTree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from instant_ppt_worker.models import DeckPlan

EDITABLE_NATIVE_SHAPE_TYPES = {
    MSO_SHAPE_TYPE.AUTO_SHAPE,
    MSO_SHAPE_TYPE.CALLOUT,
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.DIAGRAM,
    MSO_SHAPE_TYPE.FREEFORM,
    MSO_SHAPE_TYPE.IGX_GRAPHIC,
    MSO_SHAPE_TYPE.LINE,
    MSO_SHAPE_TYPE.TABLE,
}


def _normalize_text(value: str) -> str:
    return " ".join(value.split())


def _shape_metrics(
    shapes: object,
    *,
    slide_width: int,
    slide_height: int,
) -> tuple[int, list[str], int, int]:
    text_count = 0
    text_values: list[str] = []
    native_shape_count = 0
    full_slide_picture_count = 0
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            nested_count, nested_text, nested_native, nested_full_slide = _shape_metrics(
                shape.shapes,
                slide_width=slide_width,
                slide_height=slide_height,
            )
            text_count += nested_count
            text_values.extend(nested_text)
            native_shape_count += nested_native
            full_slide_picture_count += nested_full_slide
            continue
        if shape.shape_type in EDITABLE_NATIVE_SHAPE_TYPES:
            native_shape_count += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            area_ratio = (int(shape.width) * int(shape.height)) / (slide_width * slide_height)
            near_origin = (
                int(shape.left) <= slide_width * 0.02
                and int(shape.top) <= slide_height * 0.02
            )
            if near_origin and area_ratio >= 0.9:
                full_slide_picture_count += 1
        if getattr(shape, "has_text_frame", False):
            text_count += 1
            text = _normalize_text(shape.text)
            if text:
                text_values.append(text)
    return text_count, text_values, native_shape_count, full_slide_picture_count


def _relationship_owner(relationship_part: str) -> str | None:
    if relationship_part == "_rels/.rels":
        return ""
    if "/_rels/" not in relationship_part or not relationship_part.endswith(".rels"):
        return None
    prefix, filename = relationship_part.rsplit("/_rels/", 1)
    return f"{prefix}/{filename[:-5]}"


def _resolve_relationship_target(owner: str, target: str) -> str | None:
    target_path = unquote(urlsplit(target).path).replace("\\", "/")
    if not target_path:
        return None
    if target_path.startswith("/"):
        resolved = posixpath.normpath(target_path.lstrip("/"))
    else:
        resolved = posixpath.normpath(posixpath.join(posixpath.dirname(owner), target_path))
    if resolved == ".." or resolved.startswith("../"):
        return None
    return resolved


def inspect_pptx(path: Path, deck: DeckPlan) -> dict[str, object]:
    findings: list[dict[str, str]] = []
    external_relationships: list[str] = []
    missing_relationship_targets: list[str] = []
    escaped_relationship_targets: list[str] = []
    media_references: set[str] = set()
    relationship_count = 0
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
        required = {
            "[Content_Types].xml",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
        }
        for missing in sorted(required - names):
            findings.append({"code": "PPTX_PART_MISSING", "severity": "sev1", "message": missing})
        slide_parts = sorted(
            name for name in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        )
        for name in sorted(item for item in names if item.endswith(".rels")):
            try:
                root = ElementTree.fromstring(archive.read(name))
            except Exception:
                findings.append(
                    {"code": "PPTX_RELATIONSHIP_XML_INVALID", "severity": "sev1", "message": name}
                )
                continue
            owner = _relationship_owner(name)
            for rel in root.iter():
                target = rel.attrib.get("Target")
                if target is None:
                    continue
                relationship_count += 1
                if rel.attrib.get("TargetMode", "").lower() == "external":
                    external_relationships.append(f"{name}:{target}")
                    continue
                if owner is None:
                    escaped_relationship_targets.append(f"{name}:{target}")
                    continue
                resolved = _resolve_relationship_target(owner, target)
                if resolved is None:
                    escaped_relationship_targets.append(f"{name}:{target}")
                    continue
                if resolved not in names:
                    missing_relationship_targets.append(f"{name}:{target} -> {resolved}")
                if resolved.startswith("ppt/media/"):
                    media_references.add(resolved)
        media = sorted(
            name for name in names if name.startswith("ppt/media/") and not name.endswith("/")
        )
        unreferenced_media = sorted(set(media) - media_references)

    presentation = Presentation(path)
    expected_titles = [slide.title for slide in sorted(deck.slides, key=lambda item: item.order)]
    editable_shapes = 0
    editable_native_shapes = 0
    full_slide_pictures = 0
    expected_text_count = 0
    matched_text_count = 0
    missing_text: list[str] = []
    ordered_slides = sorted(deck.slides, key=lambda item: item.order)
    slide_pairs = zip(presentation.slides, ordered_slides, strict=False)
    for index, (slide, slide_plan) in enumerate(slide_pairs):
        slide_text_count, slide_text, native_count, full_slide_count = _shape_metrics(
            slide.shapes,
            slide_width=int(presentation.slide_width),
            slide_height=int(presentation.slide_height),
        )
        editable_shapes += slide_text_count
        editable_native_shapes += native_count
        full_slide_pictures += full_slide_count
        expected = Counter(_normalize_text(value) for value in [slide_plan.title, *slide_plan.body])
        actual = Counter(slide_text)
        expected_text_count += sum(expected.values())
        for value, required_count in expected.items():
            matched = min(required_count, actual[value])
            matched_text_count += matched
            missing_text.extend(
                f"slide {index + 1}: {value}" for _ in range(required_count - matched)
            )
    if len(slide_parts) != len(expected_titles) or len(presentation.slides) != len(expected_titles):
        findings.append(
            {
                "code": "PPTX_SLIDE_COUNT_MISMATCH",
                "severity": "sev1",
                "message": (
                    f"expected {len(expected_titles)}, package {len(slide_parts)}, "
                    f"API {len(presentation.slides)}"
                ),
            }
        )
    if missing_text:
        findings.append(
            {
                "code": "PPTX_EDITABLE_TEXT_MISSING",
                "severity": "sev1",
                "message": "; ".join(missing_text),
            }
        )
    if editable_native_shapes < len(expected_titles):
        findings.append(
            {
                "code": "PPTX_EDITABLE_NATIVE_SHAPE_COVERAGE_LOW",
                "severity": "sev1",
                "message": (
                    f"expected at least {len(expected_titles)} editable native shapes, "
                    f"found {editable_native_shapes}"
                ),
            }
        )
    if full_slide_pictures:
        findings.append(
            {
                "code": "PPTX_FULL_SLIDE_BITMAP_FALLBACK",
                "severity": "sev2",
                "message": f"found {full_slide_pictures} full-slide picture fallback(s)",
            }
        )
    if external_relationships:
        findings.append(
            {
                "code": "PPTX_EXTERNAL_RELATIONSHIP",
                "severity": "sev1",
                "message": "; ".join(external_relationships),
            }
        )
    if escaped_relationship_targets:
        findings.append(
            {
                "code": "PPTX_RELATIONSHIP_TARGET_ESCAPES_PACKAGE",
                "severity": "sev1",
                "message": "; ".join(escaped_relationship_targets),
            }
        )
    if missing_relationship_targets:
        findings.append(
            {
                "code": "PPTX_RELATIONSHIP_TARGET_MISSING",
                "severity": "sev1",
                "message": "; ".join(missing_relationship_targets),
            }
        )
    if unreferenced_media:
        findings.append(
            {
                "code": "PPTX_MEDIA_UNREFERENCED",
                "severity": "sev2",
                "message": "; ".join(unreferenced_media),
            }
        )
    return {
        "schema": "instant-ppt.pptx-package-qa.v1",
        "passed": not any(item["severity"] in {"sev1", "sev2"} for item in findings),
        "slideCount": len(presentation.slides),
        "editableTextShapeCount": editable_shapes,
        "expectedEditableTextCount": expected_text_count,
        "matchedEditableTextCount": matched_text_count,
        "editableNativeShapeCount": editable_native_shapes,
        "fullSlidePictureCount": full_slide_pictures,
        "relationshipCount": relationship_count,
        "mediaParts": media,
        "mediaReferences": sorted(media_references),
        "unreferencedMediaParts": unreferenced_media,
        "missingRelationshipTargets": missing_relationship_targets,
        "externalRelationships": external_relationships,
        "findings": findings,
    }


def write_package_report(path: Path, report: dict[str, object]) -> None:
    path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
