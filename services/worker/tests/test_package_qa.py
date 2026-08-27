from __future__ import annotations

import zipfile
from collections.abc import Callable
from pathlib import Path

from instant_ppt_worker.models import DeckPlan
from instant_ppt_worker.package_qa import inspect_pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


def _deck(*, body: str = "Every expected text remains editable") -> DeckPlan:
    return DeckPlan.model_validate(
        {
            "schemaVersion": 1,
            "snapshotId": "01ARZ3NDEKTSV4RRFFQ69G5FAA",
            "title": "Package QA",
            "modeId": "native",
            "templateBinding": {
                "schemaVersion": 1,
                "templateId": "01ARZ3NDEKTSV4RRFFQ69G5FAB",
                "templateVersionId": "01ARZ3NDEKTSV4RRFFQ69G5FAC",
                "compatibilityVersion": "test",
                "roleBindings": {"content": "blank"},
            },
            "slides": [
                {
                    "schemaVersion": 1,
                    "slideId": "01ARZ3NDEKTSV4RRFFQ69G5FAD",
                    "outlineSlideId": "01ARZ3NDEKTSV4RRFFQ69G5FAE",
                    "order": 0,
                    "role": "content",
                    "title": "Relationship integrity",
                    "body": [body],
                    "editable": True,
                }
            ],
        }
    )


def _presentation(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(0.5))
    title.text = "Relationship integrity"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.5))
    body.text = "Every expected text remains editable"
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(3), Inches(1))
    presentation.save(path)


def _replace_zip_member(path: Path, member: str, transform: Callable[[bytes], bytes]) -> None:
    with zipfile.ZipFile(path) as source:
        members = [(info, source.read(info)) for info in source.infolist()]
    with zipfile.ZipFile(path, "w") as target:
        for info, data in members:
            if info.filename == member:
                data = transform(data)
            target.writestr(info, data)


def test_package_qa_proves_text_shapes_and_relationship_integrity(tmp_path: Path) -> None:
    path = tmp_path / "valid.pptx"
    _presentation(path)

    report = inspect_pptx(path, _deck())

    assert report["passed"] is True
    assert report["expectedEditableTextCount"] == 2
    assert report["matchedEditableTextCount"] == 2
    assert report["editableNativeShapeCount"] >= 1
    assert report["relationshipCount"] > 0
    assert report["missingRelationshipTargets"] == []
    assert report["unreferencedMediaParts"] == []


def test_package_qa_can_verify_agent_contract_text_instead_of_planning_copy(
    tmp_path: Path,
) -> None:
    path = tmp_path / "agent-authored.pptx"
    _presentation(path)
    deck = _deck(body="Blueprint planning copy may be condensed by the executor")

    report = inspect_pptx(
        path,
        deck,
        expected_editable_text={deck.slides[0].slide_id: [deck.slides[0].title]},
    )

    assert report["passed"] is True, report
    assert report["expectedEditableTextCount"] == 1
    assert report["matchedEditableTextCount"] == 1


def test_package_qa_accepts_title_split_across_adjacent_editable_shapes(
    tmp_path: Path,
) -> None:
    path = tmp_path / "split-title.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
    first.text = "Relationship"
    second = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(4), Inches(0.5))
    second.text = " integrity"
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(3), Inches(1))
    presentation.save(path)
    deck = _deck()

    report = inspect_pptx(
        path,
        deck,
        expected_editable_text={deck.slides[0].slide_id: [deck.slides[0].title]},
    )

    assert report["passed"] is True, report
    assert report["matchedEditableTextCount"] == 1


def test_package_qa_matches_markdown_escape_to_visible_punctuation(tmp_path: Path) -> None:
    path = tmp_path / "markdown-escape.pptx"
    _presentation(path)

    report = inspect_pptx(path, _deck(body=r"Every expected text remains edit\-able"))

    assert report["passed"] is False
    assert any(
        finding["code"] == "PPTX_EDITABLE_TEXT_MISSING"
        for finding in report["findings"]
    )

    # The same comparison succeeds when the PPTX contains the rendered hyphen.
    presentation = Presentation(path)
    presentation.slides[0].shapes[1].text = "Every expected text remains edit-able"
    presentation.save(path)
    report = inspect_pptx(path, _deck(body=r"Every expected text remains edit\-able"))
    assert report["passed"] is True


def test_package_qa_rejects_missing_media_relationship_target(tmp_path: Path) -> None:
    path = tmp_path / "broken.pptx"
    _presentation(path)
    relationship = (
        b'<Relationship Id="rId999" '
        b'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" '
        b'Target="../media/missing.png"/>'
    )

    def add_missing_media(data: bytes) -> bytes:
        return data.replace(b"</Relationships>", relationship + b"</Relationships>")

    _replace_zip_member(path, "ppt/slides/_rels/slide1.xml.rels", add_missing_media)
    report = inspect_pptx(path, _deck())

    assert report["passed"] is False
    assert report["mediaReferences"] == ["ppt/media/missing.png"]
    assert report["missingRelationshipTargets"]
    assert any(
        finding["code"] == "PPTX_RELATIONSHIP_TARGET_MISSING"
        for finding in report["findings"]
    )
